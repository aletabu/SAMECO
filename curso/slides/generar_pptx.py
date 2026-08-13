# Genera los 4 PPT del curso SAMECO a partir del guion de 00_Plan_del_Curso.md.
# Uso: python3 generar_pptx.py   (crea/reescribe los .pptx en esta carpeta)
# Audiencia no técnica: poco texto por lámina; el detalle va en las notas del orador.

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK = RGBColor(0x1F, 0x29, 0x37)      # texto principal
GRIS = RGBColor(0x6B, 0x72, 0x80)     # texto secundario
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF3, 0xF4, 0xF6)    # fondo de paneles

ACENTOS = {
    1: RGBColor(0x25, 0x63, 0xEB),  # azul
    2: RGBColor(0x05, 0x96, 0x69),  # verde
    3: RGBColor(0xD9, 0x77, 0x06),  # ámbar
    4: RGBColor(0x7C, 0x3A, 0xED),  # violeta
}

ANCHO, ALTO = Inches(13.333), Inches(7.5)
FUENTE = "Arial"
LOGO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "logo-sameco-trim.png")
LOGO_RATIO = 440 / 83  # ancho/alto del PNG recortado


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _texto(tf, texto, size, color, bold=False, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = texto
    r.font.name = FUENTE
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return p


class Deck:
    def __init__(self, sesion, archivo):
        self.n = sesion
        self.acento = ACENTOS[sesion]
        self.archivo = archivo
        self.prs = Presentation()
        self.prs.slide_width = ANCHO
        self.prs.slide_height = ALTO
        self.blank = self.prs.slide_layouts[6]

    def _slide(self, notas=None):
        s = self.prs.slides.add_slide(self.blank)
        if notas:
            s.notes_slide.notes_text_frame.text = notas
        return s

    def _pie(self, s, texto):
        box = s.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4))
        _texto(box.text_frame, texto, 10, GRIS)

    def _logo(self, s, alto_in, x_in, y_in):
        s.shapes.add_picture(LOGO, Inches(x_in), Inches(y_in),
                             width=Inches(alto_in * LOGO_RATIO), height=Inches(alto_in))

    def _cabecera(self, s, titulo, kicker=None):
        barra = s.shapes.add_shape(1, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.9))
        _fill(barra, self.acento)
        self._logo(s, 0.3, 11.05, 0.55)
        if kicker:
            kb = s.shapes.add_textbox(Inches(0.95), Inches(0.42), Inches(11.5), Inches(0.4))
            _texto(kb.text_frame, kicker.upper(), 12, self.acento, bold=True)
            ty = Inches(0.78)
        else:
            ty = Inches(0.62)
        tb = s.shapes.add_textbox(Inches(0.95), ty, Inches(11.7), Inches(0.9))
        _texto(tb.text_frame, titulo, 30, INK, bold=True)

    # ---- tipos de lámina ----

    def portada(self, titulo, subtitulo, notas=None):
        s = self._slide(notas)
        fondo = s.shapes.add_shape(1, 0, 0, ANCHO, ALTO)
        _fill(fondo, INK)
        franja = s.shapes.add_shape(1, 0, Inches(4.9), ANCHO, Inches(0.12))
        _fill(franja, self.acento)
        self._logo(s, 0.5, 0.9, 0.75)
        kb = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5))
        _texto(kb.text_frame, f"ASISTENTE Q&A SAMECO — SESIÓN {self.n} DE 4", 15, self.acento, bold=True)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(2.25), Inches(11.5), Inches(1.8))
        _texto(tb.text_frame, titulo, 44, BLANCO, bold=True)
        sb = s.shapes.add_textbox(Inches(0.9), Inches(5.25), Inches(11.5), Inches(0.9))
        _texto(sb.text_frame, subtitulo, 17, RGBColor(0xD1, 0xD5, 0xDB))
        return s

    def seccion(self, titulo, bajada=None, notas=None):
        s = self._slide(notas)
        fondo = s.shapes.add_shape(1, 0, 0, ANCHO, ALTO)
        _fill(fondo, self.acento)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.4))
        _texto(tb.text_frame, titulo, 40, BLANCO, bold=True)
        if bajada:
            bb = s.shapes.add_textbox(Inches(0.9), Inches(4.25), Inches(11.5), Inches(1.0))
            _texto(bb.text_frame, bajada, 18, BLANCO)
        return s

    def bullets(self, titulo, items, kicker=None, notas=None, pie=None):
        """items: lista de str o (str, nivel). Nivel 1 = sub-bullet."""
        s = self._slide(notas)
        self._cabecera(s, titulo, kicker)
        box = s.shapes.add_textbox(Inches(0.95), Inches(1.9), Inches(11.5), Inches(4.9))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for it in items:
            txt, nivel = it if isinstance(it, tuple) else (it, 0)
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = nivel
            p.space_after = Pt(14 if nivel == 0 else 8)
            r = p.add_run()
            r.text = ("•  " if nivel == 0 else "–  ") + txt
            r.font.name = FUENTE
            r.font.size = Pt(20 if nivel == 0 else 16)
            r.font.color.rgb = INK if nivel == 0 else GRIS
        if pie:
            self._pie(s, pie)
        return s

    def concepto(self, termino, definicion, analogia, notas=None):
        """Lámina de concepto: término grande + definición + analogía de biblioteca."""
        s = self._slide(notas)
        kb = s.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.4))
        _texto(kb.text_frame, "CONCEPTO", 13, self.acento, bold=True)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(1.25), Inches(11.5), Inches(1.2))
        _texto(tb.text_frame, termino, 48, INK, bold=True)
        db = s.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.6))
        _texto(db.text_frame, definicion, 24, INK)
        panel = s.shapes.add_shape(1, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.5))
        _fill(panel, PANEL)
        panel.text_frame.margin_left = Inches(0.35)
        panel.text_frame.margin_top = Inches(0.22)
        panel.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        _texto(panel.text_frame, "En la biblioteca:  " + analogia, 17, GRIS)
        return s

    def panel(self, titulo, lineas, kicker=None, notas=None):
        """Bloque de texto citado sobre fondo gris (p. ej. instrucciones de ejemplo)."""
        s = self._slide(notas)
        self._cabecera(s, titulo, kicker)
        panel = s.shapes.add_shape(1, Inches(0.95), Inches(1.95), Inches(11.45), Inches(4.6))
        _fill(panel, PANEL)
        tf = panel.text_frame
        tf.margin_left = Inches(0.4)
        tf.margin_top = Inches(0.35)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        first = True
        for ln in lineas:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(10)
            r = p.add_run()
            r.text = ln
            r.font.name = FUENTE
            r.font.size = Pt(16)
            r.font.color.rgb = INK
            r.font.bold = ln.isupper()
        return s

    def tabla(self, titulo, headers, filas, kicker=None, notas=None, anchos=None, pie=None):
        s = self._slide(notas)
        self._cabecera(s, titulo, kicker)
        ncols = len(headers)
        alto_tabla = Inches(min(0.55 * (len(filas) + 1) + 0.3, 4.9))
        gt = s.shapes.add_table(len(filas) + 1, ncols, Inches(0.95), Inches(1.95),
                                Inches(11.45), alto_tabla)
        t = gt.table
        if anchos:
            for i, a in enumerate(anchos):
                t.columns[i].width = Inches(a)
        for j, h in enumerate(headers):
            c = t.cell(0, j)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = self.acento
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FUENTE
                    r.font.size = Pt(14)
                    r.font.bold = True
                    r.font.color.rgb = BLANCO
        for i, fila in enumerate(filas, start=1):
            for j, val in enumerate(fila):
                c = t.cell(i, j)
                c.text = str(val)
                c.fill.solid()
                c.fill.fore_color.rgb = BLANCO if i % 2 else PANEL
                for p in c.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = FUENTE
                        r.font.size = Pt(13)
                        r.font.color.rgb = INK
        if pie:
            self._pie(s, pie)
        return s

    def guardar(self):
        self.prs.save(self.archivo)
        print("OK", self.archivo)


# ================= SESIÓN 1 — Conceptos y descubrimiento =================

d = Deck(1, "Sesion_1_Conceptos_y_Descubrimiento.pptx")

d.portada(
    "Un asistente que responde\npor SAMECO",
    "Sesión 1 — Conceptos y descubrimiento  ·  Qué vamos a construir y con qué materiales",
    notas="Antes de esta lámina: demo en vivo del sandbox YA funcionando (0-10 min). Tres preguntas: "
          "una de evento, una de histórico, una que no está (fallback). Ver el final antes de empezar "
          "motiva más que cualquier teoría.")

d.tabla("Agenda de hoy", ["", "Bloque"],
        [["0–10", "Demo: el asistente funcionando"],
         ["10–35", "Los conceptos con los que vamos a trabajar"],
         ["35–50", "El diseño: un asistente, dos bibliotecas"],
         ["50–80", "Taller: inventario de contenidos"],
         ["80–90", "Cierre y tareas"]],
        anchos=[1.5, 9.95],
        notas="Respaldo: Guía Técnica Parte C + Guía de Capacitación §1.")

d.bullets("Lo que acaban de ver", [
    "Un asistente que responde preguntas sobre el encuentro y sobre trabajos históricos",
    "Siempre dice de qué documento sacó la respuesta",
    "Cuando no sabe, lo dice — no inventa",
    "Eso mismo, con los documentos reales de SAMECO, es lo que vamos a construir juntos",
], kicker="Demo",
   notas="Recalcar: lo que vieron corre sobre material de práctica ficticio. El objetivo del curso es "
         "que al final el asistente corra sobre SU material y lo operen ustedes.")

d.seccion("Los conceptos", "Cinco ideas — con una analogía que nos va a acompañar todo el curso",
          notas="Regla didáctica: cero código. Todo se ancla en la analogía de la biblioteca.")

d.concepto("IA generativa (LLM)",
           "Un programa que aprendió a redactar leyendo millones de textos. Escribe muy bien… "
           "pero no conoce los documentos de SAMECO.",
           "un redactor brillante recién llegado a la ciudad: escribe perfecto, no sabe nada del barrio.",
           notas="No entrar en detalle técnico. La idea clave: sabe REDACTAR, no sabe QUÉ PASÓ en SAMECO. "
                 "Eso prepara el problema de la alucinación.")

d.concepto("Alucinación",
           "Cuando no sabe la respuesta, la inventa — con total seguridad y muy bien escrita.",
           "un redactor que jamás dice \"no sé\": si le preguntás por un libro que no existe, te lo resume igual.",
           notas="MOMENTO CLAVE de la sesión: mostrar UNA alucinación real en vivo — preguntarle a un chat "
                 "genérico por la agenda del encuentro SAMECO 2026 y ver cómo inventa horarios y oradores. "
                 "Es el momento 'ahh' que justifica todo lo demás.")

d.concepto("RAG: buscar primero, responder después",
           "Antes de responder, el asistente busca en los documentos cargados y redacta SOLO "
           "a partir de lo que encontró.",
           "el bibliotecario no opina de memoria: va al estante, trae el libro y te lee lo que dice.",
           notas="RAG = Retrieval-Augmented Generation, pero la sigla no importa; importa el orden: "
                 "primero buscar, después redactar. Esta es la diferencia entre nuestro asistente y un chat genérico.")

d.concepto("Grounding y citas",
           "El asistente está obligado a basarse en los documentos, y muestra la fuente de cada "
           "respuesta. Si no está en los documentos: \"no cuento con esa información\".",
           "cada respuesta viene con la ficha del libro del que salió. Y si el libro no está, te lo dice.",
           notas="Conectar con la demo: mostrar dónde aparecen las citas en la respuesta. El fallback honesto "
                 "es una FEATURE, no una falla — diferenciador clave frente al chat genérico de la lámina anterior.")

d.concepto("Contexto",
           "El asistente recuerda la conversación: se puede repreguntar \"¿y el segundo día?\" "
           "sin repetir todo.",
           "el bibliotecario se acuerda de qué venían hablando — no hace falta arrancar de cero en cada pregunta.",
           notas="Demo rápida si hay tiempo: pregunta 8 del set (¿y el segundo día a la mañana?). "
                 "También mencionar el límite: el contexto dura la conversación, no entre visitas.")

d.tabla("La biblioteca: el diccionario del curso",
        ["Cuando escuchen…", "Piensen en…"],
        [["Datastore", "La biblioteca (los estantes con los documentos)"],
         ["Indexar", "Fichar los libros para poder encontrarlos rápido"],
         ["El agente / la app", "El bibliotecario que atiende la ventanilla"],
         ["Prompt / instrucciones", "Las reglas de trabajo del bibliotecario"],
         ["Bucket (Cloud Storage)", "El estante público: lo que el asistente puede leer"],
         ["Widget", "La ventanilla de consultas puesta en la página web"],
         ["OCR", "Transcribir un documento viejo escaneado para poder ficharlo"]],
        kicker="Glosario", anchos=[3.4, 8.05],
        notas="Esta lámina queda como referencia permanente: son los términos que van a escuchar "
              "durante todo el curso. Sugerencia: imprimirla o dejarla fijada en el chat del equipo.")

d.bullets("El diseño: un asistente, dos bibliotecas", [
    "Biblioteca \"Evento octubre 2026\": agenda, oradores, inscripción, logística",
    "Biblioteca \"Archivo histórico\": trabajos de mejora continua de años anteriores",
    "Un solo asistente decide en cuál buscar según la pregunta",
    ("¿Por qué no dos asistentes? Una sola ventanilla es más simple para el visitante y para ustedes", 1),
    "Dos lugares de uso: la página web del evento (público) y el equipo interno",
], kicker="Diseño",
   notas="Preguntar acá cómo tienen organizado el material en Drive (ya sabemos que tienen Workspace). "
         "Explicar con la lámina siguiente el rol de Drive vs. el estante público.")

d.bullets("Su Drive sigue siendo su lugar de trabajo", [
    "Drive = el escritorio: ahí editan, versionan y guardan todo, como siempre",
    "El estante público (bucket) = lo que el asistente puede leer y citar",
    "\"Publicar\" un documento = copiarlo del escritorio al estante (un botón, 30 segundos)",
    "Nada llega al asistente por accidente: publicar es un acto deliberado",
    ("Borradores y versiones viejas se quedan en Drive, invisibles para el público", 1),
], kicker="Diseño",
   notas="Razón técnica (no contarla salvo que pregunten): el widget público no puede servir contenido "
         "desde Drive porque los datastores de Drive llevan control de acceso. Verificado ago-2026.")

d.bullets("Taller: el inventario de contenidos", [
    "En una planilla compartida: qué documentos existen, dónde están, en qué formato",
    "Clasificamos con un semáforo:",
    ("VERDE — texto limpio (Word, PDF moderno): entra directo", 1),
    ("AMARILLO — escaneos: necesitan OCR (transcripción automática)", 1),
    ("ROJO — láminas, pósters, PPT muy visuales: hay que describirlos", 1),
    "De paso anotamos la ficha de cada uno: año, autor, tema — la vamos a usar en la Sesión 2",
    "Prioridad: lo que la gente más consulta va primero",
], kicker="Taller · 30 min",
   notas="Guiar con la planilla proyectada. Semáforo = sección 4 de la propuesta. Gestionar expectativa "
         "sobre escaneos muy deteriorados desde YA: los ilegibles van a transcripción manual o segunda tanda.")

d.bullets("Tareas para la próxima sesión", [
    "Equipo: completar el inventario y juntar los archivos en una carpeta de Drive",
    "Alejandro: crear el proyecto de SAMECO en Google Cloud y dejar la estructura lista",
    "Próxima sesión: cargamos SU material y el asistente empieza a responder con él",
], kicker="Cierre",
   notas="Entre sesiones: crear proyecto, billing con alerta, bucket con estructura, y cargar un primer "
         "lote de documentos fáciles para que la Sesión 2 arranque con algo real.")

d.guardar()

# ================= SESIÓN 2 — Preparación del material y carga =================

d = Deck(2, "Sesion_2_Preparacion_y_Carga.pptx")

d.portada("La biblioteca se llena hoy",
          "Sesión 2 — Preparación del material y carga  ·  De la carpeta de Drive al asistente",
          notas="Respaldo: propuesta §4 + Capacitación Embeddings Parte 1-2 (solo ideas, sin código). "
                "Objetivo: que sepan dejar la biblioteca lista y entiendan por qué un documento 'entra mal'.")

d.tabla("Agenda de hoy", ["", "Bloque"],
        [["0–15", "Repaso y estado del inventario"],
         ["15–40", "Por qué no todo documento entra igual"],
         ["40–70", "Carga en vivo — el mouse lo manejan ustedes"],
         ["70–85", "El procedimiento \"sumar un documento\""],
         ["85–90", "Tareas"]],
        anchos=[1.5, 9.95])

d.bullets("¿Cómo venimos con el inventario?", [
    "Repasamos la planilla: ¿qué apareció? ¿qué falta?",
    "Dudas de clasificación: ¿este PPT es verde o rojo?",
    "Elegimos el primer lote a cargar hoy: los verdes más consultados",
], kicker="Repaso · 15 min",
   notas="Resolver dudas de semáforo con ejemplos concretos de su material.")

d.tabla("Por qué no todo documento entra igual",
        ["El documento", "Qué pasa al cargarlo"],
        [["Word / PDF moderno (texto)", "Entra directo — se puede seleccionar el texto"],
         ["PowerPoint", "El texto entra; lo que es imagen (fotos de láminas), no"],
         ["Escaneo (PDF-foto)", "Necesita OCR: la transcripción automática del escaneo"],
         ["Lámina, póster, A3", "Hay que describirlo en texto para que el asistente lo \"vea\""]],
        kicker="Concepto", anchos=[3.9, 7.55],
        notas="Truco para explicarlo: '¿pueden seleccionar el texto con el mouse? Entonces el asistente "
              "también puede leerlo. ¿No pueden? El asistente tampoco.'")

d.bullets("Demo: el trabajo escaneado de 2009", [
    "Un PDF que es una foto: no se puede seleccionar ni una palabra",
    "Y sin embargo… el asistente responde qué scrap tenía la línea de pintura",
    "Eso es el OCR: la biblioteca transcribió el documento viejo al ficharlo",
    "Límite honesto: si el escaneo está muy deteriorado, la transcripción sale mal",
], kicker="Demo",
   notas="Demo con el escaneo simulado 2009 del sandbox: mostrar el PDF (texto no seleccionable) y la "
         "respuesta correcta (8,5% → 2,1%). Cerrar gestionando expectativa: los escaneos ilegibles de su "
         "archivo van a transcripción manual o segunda tanda.")

d.bullets("Fortalecer la biblioteca: la ficha de cada documento", [
    "Además del documento, la biblioteca puede guardar su ficha: título, año, autor, sector, tema y un resumen corto",
    "Con ficha, el bibliotecario encuentra mejor:",
    ("\"¿Qué trabajos hay de 2018?\" deja de depender de que el año figure en el texto", 1),
    ("\"¿Hay algo del sector alimentario?\" funciona aunque el documento nunca diga \"alimentario\"", 1),
    "Análisis previo con IA: se le hace leer cada documento y propone resumen y etiquetas — ustedes las revisan",
    "División del trabajo: el equipo completa la ficha en la planilla del inventario; Alejandro la carga a la biblioteca",
], kicker="Concepto",
   notas="Detalle técnico (no contar en clase): esto se implementa con la opción 'Documentos con metadatos "
         "(RAG)' del datastore — un JSONL que vincula metadatos a cada archivo del bucket; el esquema se "
         "detecta automático. Trabajo entre sesiones de Alejandro: script que lee cada documento, genera "
         "resumen + etiquetas con IA (análisis semántico), y arma el JSONL con lo validado en la planilla. "
         "ACLARACIÓN DE VOCABULARIO por si surge: esto NO es fine-tuning (no se reentrena ningún modelo); "
         "es enriquecimiento de datos — la palanca correcta y barata para mejorar un RAG gestionado.")

d.bullets("Carga en vivo — hoy el mouse es de ustedes", [
    "Subimos su primer lote real al estante público (bucket)",
    "Creamos / actualizamos las dos bibliotecas (la histórica, con OCR activado)",
    "Vemos la indexación en curso: el fichado tarda de minutos a una hora",
    "Regla de hoy: el teclado lo maneja una persona del equipo, no Alejandro",
], kicker="Práctica · 30 min",
   notas="Acá empieza la transferencia. Plan B obligatorio: si la indexación se demora, mostrar el "
         "sandbox propio mientras tanto. Nunca cargar por primera vez en vivo sin red de seguridad.")

d.bullets("El procedimiento \"sumar un documento\"", [
    "1.  Guardarlo en Drive como siempre (escritorio)",
    "2.  Subirlo al estante público (botón Upload en el navegador)",
    "3.  Esperar la indexación (vamos a cronometrar cuánto tarda)",
    "4.  Verificar: hacerle al asistente una pregunta que solo ese documento responde",
    "Este circuito es el corazón de su autonomía — lo van a hacer sin ayuda",
], kicker="Procedimiento",
   notas="Ensayarlo COMPLETO una vez en vivo y cronometrar: ese número responde la futura pregunta "
         "'¿cuándo aparece lo que subo?'. Este procedimiento va al manual de operación de la Sesión 4.")

d.bullets("Tareas para la próxima sesión", [
    "Equipo: terminar de cargar el material verde; separar escaneos y láminas problemáticos",
    "Equipo: completar las fichas (año, autor, tema, resumen) en la planilla del inventario",
    "Alejandro: procesar el material difícil (OCR, descripciones) y cargar las fichas a la biblioteca",
    "Próxima sesión: el asistente responde con TODO su material — y aprendemos a evaluarlo",
], kicker="Cierre",
   notas="Entre sesiones: OCR/descripciones del material difícil, completar indexación y crear la app "
         "conectada a los 2 datastores.")

d.guardar()

# ================= SESIÓN 3 — El asistente en marcha: prompt y pruebas =================

d = Deck(3, "Sesion_3_Prompt_y_Pruebas.pptx")

d.portada("Enseñarle a responder bien",
          "Sesión 3 — El asistente en marcha  ·  Instrucciones, pruebas e iteración",
          notas="Respaldo: Guía de Capacitación §2 (diseño del prompt) — EL documento de esta sesión. "
                "Objetivo: que sepan evaluar respuestas e iterar el comportamiento.")

d.tabla("Agenda de hoy", ["", "Bloque"],
        [["0–10", "El asistente ya responde con su material"],
         ["10–30", "Anatomía del comportamiento: las 4 decisiones"],
         ["30–70", "Taller: su set de preguntas de prueba"],
         ["70–85", "Iterar: un cambio por vez"],
         ["85–90", "Tareas"]],
        anchos=[1.5, 9.95])

d.bullets("El asistente ya responde", [
    "Preview con el material real de SAMECO cargado en la Sesión 2",
    "Hoy no construimos: hoy evaluamos y afinamos",
    "La pregunta de la sesión: ¿cómo sabemos si responde BIEN?",
], kicker="Demo · 10 min")

d.bullets("Las 4 decisiones que definen el comportamiento", [
    "ROL — quién es: \"el asistente oficial de SAMECO\"",
    "TONO — cómo habla: breve, cordial, en español",
    "REGLAS — qué puede y qué no: solo responder con los documentos; si no está, decirlo",
    "CONTEXTO — qué recuerda: la conversación en curso",
    "Todo eso se escribe en las instrucciones del bibliotecario (el prompt)",
], kicker="Concepto",
   notas="Mostrar el con/sin instrucciones ensayado en el lab (paso 7): la MISMA pregunta con dos "
         "comportamientos distintos. Ese contraste explica más que cualquier teoría.")

d.bullets("Misma pregunta, dos bibliotecarios", [
    "Sin instrucciones: respuesta larga, genérica, a veces inventa",
    "Con instrucciones: breve, cita la fuente, y si no sabe lo dice",
    "La diferencia no es el motor: son las reglas de trabajo",
], kicker="Demo",
   notas="Demo en vivo del contraste con/sin preamble. Técnico (no contar): el prompt del sistema se "
         "llama preamble y el ajuste fino es por API; desde consola el control es limitado.")

d.bullets("Las 3 reglas de unas buenas instrucciones", [
    "CUÁNDO BUSCAR — siempre en la biblioteca, nunca \"de memoria\"",
    ("La regla se escribe en positivo y en imperativo: \"buscá SIEMPRE antes de responder\"", 1),
    "CÓMO CITAR — cada dato con su fuente, en un formato fijo",
    ("Si la cita es opcional, desaparece; si el formato es libre, es inconsistente", 1),
    "QUÉ HACER SI NO ESTÁ — el fallback se define por adelantado, con derivación",
    ("No alcanza con \"no lo sé\": se indica a dónde acudir (mail de la organización)", 1),
], kicker="Concepto",
   notas="Estas tres reglas son el esqueleto de cualquier prompt de RAG bien hecho. Si falta una, el "
         "agujero aparece en las pruebas: sin regla de búsqueda → responde de memoria; sin formato de "
         "cita → citas a veces sí a veces no; sin fallback definido → inventa cuando no encuentra.")

d.panel("Ejemplo: las instrucciones del asistente SAMECO", [
    "Sos el asistente oficial del Encuentro SAMECO 2026.",
    "REGLAS:",
    "1.  Buscá SIEMPRE en los documentos antes de responder. Nunca respondas con "
    "conocimiento propio ni de memoria.",
    "2.  Cada dato de la respuesta debe indicar de qué documento salió.",
    "3.  Si la información no está en los documentos, respondé: \"No cuento con esa "
    "información\" y sugerí escribir a la organización.",
    "4.  Respondé en español, breve y cordial.",
], kicker="Ejemplo",
   notas="Leerlo en voz alta y mapear cada línea con las 3 reglas de la lámina anterior (la 4 es tono). "
         "Técnico (no contar): esto va en el preamble (promptSpec.preamble de la Answer API); desde la "
         "consola el control es limitado. Confirmar con SAMECO el mail real de derivación antes de fijarlo.")

d.tabla("Los 4 tipos de pregunta que el set tiene que cubrir",
        ["Tipo", "Ejemplo", "Qué esperamos"],
        [["Factual directa", "\"¿Cuánto sale la inscripción?\"", "Respuesta puntual con su cita"],
         ["De síntesis", "\"¿Qué factor de éxito se repite en los trabajos?\"",
          "Combina varios documentos, cita varios"],
         ["Fuera de la biblioteca", "\"¿Va a estar Carlos Pagni?\"",
          "Fallback honesto + derivación — nunca inventar"],
         ["Ambigua", "\"¿Y los beneficios?\" / \"contame más\"",
          "Repregunta, o responde lo más relevante sin divagar"]],
        kicker="Método", anchos=[2.6, 4.9, 3.95],
        notas="Taxonomía estándar de evaluación de asistentes RAG. El set 03 ya cubre las tres primeras "
              "(las 'canarias' 5/10/16/20 son las de fuera de la biblioteca); las ambiguas conviene "
              "agregarlas como bloque nuevo al armar el set de SAMECO. La pregunta 8 ('¿y el segundo "
              "día?') es un caso de ambigua que se resuelve por contexto conversacional.")

d.bullets("Taller: su set de preguntas de prueba", [
    "Armamos 15–20 preguntas fijas, en tres bloques:",
    ("A — Evento: fechas, precios, horarios, logística", 1),
    ("B — Histórico: trabajos, autores, resultados", 1),
    ("C — Mixtas: preguntas que cruzan o confunden los dos dominios", 1),
    "Cuidamos que estén los 4 tipos: factuales, de síntesis, canarias y ambiguas",
    "Las corremos en vivo y anotamos el resultado de cada una",
], kicker="Taller · 40 min",
   notas="Plantilla: 03_Set_Preguntas_Prueba.md (bloques A/B/C). Las canarias detectan alucinación: "
         "si el asistente inventa ahí, el grounding está mal y es lo primero a arreglar.")

d.bullets("¿Qué hace buena a una respuesta?", [
    "FUNDAMENTADA — la información sale de los documentos, no de la \"memoria\" del modelo",
    "CITADA — cada afirmación concreta tiene su fuente",
    "RELEVANTE — responde LO que se preguntó, sin relleno",
    "COMPLETA — trae toda la información necesaria, no un fragmento",
    "Una respuesta puede estar bien citada y aun así ser incompleta — por eso se evalúan las cuatro",
], kicker="Método",
   notas="Los 4 criterios estándar de calidad RAG: grounding, citation, relevance, completeness. "
         "Conectar con la lámina siguiente: la escala de anotación es la versión operativa de esto.")

d.bullets("Cómo anotar cada respuesta", [
    "CORRECTA — respondió bien y citó la fuente",
    "INCOMPLETA — le faltó parte de la respuesta",
    "INVENTADA — dijo algo que no está en los documentos (¡la peor!)",
    "FALLBACK CORRECTO — no estaba, y lo dijo honestamente",
    "La planilla con estas notas es la memoria de calidad del asistente",
], kicker="Método",
   notas="Planilla: fecha, versión de instrucciones, resultado por pregunta. Sirve para comparar "
         "versiones de forma pareja — y es el registro que Alejandro usa entre sesiones.")

d.bullets("Iterar: un cambio por vez", [
    "Elegimos la peor falla de la corrida",
    "Cambiamos UNA sola cosa: una instrucción, o un documento fuente",
    "Volvemos a correr el set completo y comparamos",
    "Regla de oro: si cambiás dos cosas y mejora, no sabés cuál fue",
], kicker="Método",
   notas="Hacer el ciclo completo UNA vez en vivo para que lo vean de punta a punta.")

d.bullets("Tareas para la próxima sesión", [
    "Equipo: correr el set completo durante la semana y anotar las fallas en la planilla",
    "Alejandro: ajustar instrucciones y fuentes según la planilla",
    "Próxima sesión: ¡publicamos! El asistente sale a la página del evento",
], kicker="Cierre",
   notas="Entre sesiones: afinar prompt y fuentes, dejar calidad estable, preparar el canal de "
         "publicación (allowlist del dominio del sitio de SAMECO).")

d.guardar()

# ================= SESIÓN 4 — Publicación, entrega y por dentro =================

d = Deck(4, "Sesion_4_Publicacion_y_Entrega.pptx")

d.portada("El asistente sale al mundo",
          "Sesión 4 — Publicación, entrega y \"cómo funciona por dentro\"",
          notas="Respaldo: propuesta §6 (entregables) + Capacitación Embeddings (demo opcional). "
                "Objetivo: asistente publicado y equipo autónomo.")

d.tabla("Agenda de hoy", ["", "Bloque"],
        [["0–20", "Publicación en vivo"],
         ["20–45", "El manual de operación"],
         ["45–60", "Costos y monitoreo"],
         ["60–80", "Opcional: cómo funciona por dentro"],
         ["80–90", "Cierre y entrega"]],
        anchos=[1.5, 9.95])

d.bullets("Publicación en vivo", [
    "Pegamos el código del widget (la ventanilla) en la página del evento",
    "Es un bloque de texto que el webmaster inserta una sola vez",
    "La primera consulta pública la hace el equipo — momento ceremonial",
], kicker="En vivo · 20 min",
   notas="Si el sitio real no está disponible, usar una página de prueba del dominio real (el dominio "
         "debe estar allowlisteado). Nota: el estilo del widget no se personaliza por CSS — decirlo "
         "antes de que lo pregunten.")

d.bullets("El manual de operación: su asistente, sus reglas", [
    "Sumar un documento: el circuito de la Sesión 2 (Drive → estante → indexar → verificar)",
    "Si una respuesta está mal, preguntarse en orden:",
    ("¿Falta el documento? → subirlo", 1),
    ("¿Está mal escaneado? → revisarlo o transcribirlo", 1),
    ("¿Hay que ajustar una instrucción? → anotar el caso y avisar a Alejandro", 1),
    "Cambios de agenda del evento: actualizar el documento, volver a publicar, verificar",
], kicker="Entregable",
   notas="Recorrer el documento de operación entregado. Dejar claro el criterio de escalamiento: "
         "contenido lo resuelven ellos; comportamiento (instrucciones) se anota y se consulta.")

d.bullets("Después del estreno: escuchar las preguntas reales", [
    "La gente pregunta distinto de lo que ensayamos — las preguntas reales son la mejor fuente de mejora",
    "La rutina semanal del equipo: mirar cuatro señales",
    ("¿Las respuestas siguen citando su fuente?", 1),
    ("¿Cuántas veces dijo \"no cuento con esa información\"? Un fallback repetido = falta un documento", 1),
    ("¿Qué temas pregunta más la gente? ¿Coinciden con lo que priorizamos?", 1),
    ("De las respuestas flojas: ¿el documento citado realmente contenía la respuesta?", 1),
    "Las mejores preguntas reales se suman al set de prueba — el set crece con el uso",
], kicker="Operación",
   notas="Idea fuerza: un fallback frecuente no es una falla del asistente, es el asistente diciéndoles "
         "qué documento falta en la biblioteca. Cierra el ciclo: preguntas reales → planilla → contenido "
         "o instrucción → set de prueba. Técnico A VERIFICAR antes de la sesión: qué métricas expone hoy "
         "la consola de AI Applications para la app (volumen, top queries) y desde dónde se ven los "
         "fallbacks; si la consola queda corta, la rutina se hace muestreando conversaciones o con la "
         "planilla de casos reportados.")

d.bullets("Costos: qué esperar cada mes", [
    "Las primeras 10.000 consultas del mes: gratis (permanente)",
    "Matiz: las \"respuestas avanzadas\" se facturan aparte (~$4 por cada 1.000)",
    "El índice de documentos: gratis hasta 10 GiB — su biblioteca entra holgada",
    "Ya quedó configurada una alerta de presupuesto: si algo se dispara, avisa solo",
    "Dónde mirarlo: sección Facturación (Billing) de la consola",
], kicker="Operación",
   notas="Cifras validadas ago-2026 (04_Actualizacion §2). Estimación realista para su tráfico: "
         "$0-40/mes. Mostrar Billing en vivo y dónde está la alerta.")

d.bullets("Opcional: ¿cómo funciona por dentro?", [
    "El documento se corta en fragmentos (fichas)",
    "Cada fragmento se convierte en un \"código de significado\" (embedding)",
    "Una pregunta busca los fragmentos más parecidos en significado, no en palabras",
    "Con esos fragmentos, el modelo redacta la respuesta y cita la fuente",
    "Todo esto la plataforma lo hace sola — por eso no programamos nada en 4 sesiones",
], kicker="Demo opcional · 15 min",
   notas="Solo si el grupo tiene curiosidad. Mini-demo RAG de la Guía Técnica B.5 (chunking → "
         "embeddings → retrieval) como espectáculo, no como clase.")

d.bullets("Cierre: qué queda entregado", [
    "El asistente publicado en la página del evento, respondiendo con su material",
    "El manual de operación: sumar, verificar, diagnosticar",
    "El set de preguntas de prueba y su planilla de calidad",
    "La alerta de costos configurada",
    "El canal de soporte con Alejandro para lo que escale",
], kicker="Entrega",
   notas="Checklist de entrega = Guía de Capacitación §4: conceptos ✓ prompt ✓ set de pruebas ✓ "
         "costos ✓. Acordar canal y tiempos de soporte post-entrega. El caché de respuestas NO se "
         "promete: es optimización post-evento si el tráfico lo justifica.")

d.guardar()
